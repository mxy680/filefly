from abc import ABC, abstractmethod
from typing import Tuple, List
from pptx import Presentation
from PIL import Image
import io
from odf.opendocument import load as odf_load
from odf.text import P
from odf.draw import Frame, Image as ODFImage
from zipfile import ZipFile


class PresentationExtractor(ABC):
    """
    Abstract base class for slideshow extraction.
    """

    @property
    def file_type(self) -> str:
        return "Slideshow"

    @abstractmethod
    def extract(self, buffer: bytes) -> Tuple[str, List[bytes]]:
        """
        Abstract method to extract text and images from a slideshow.

        Args:
            buffer (bytes): The file content.

        Returns:
            Tuple[str, List[bytes]]: Extracted text and images.
        """
        pass


class PPTXExtractor(PresentationExtractor):
    """
    Extracts text and images from PPTX files.
    """

    def extract(self, buffer: bytes) -> Tuple[str, List[bytes]]:
        try:
            # Open the PPTX file
            presentation = Presentation(io.BytesIO(buffer))
            text = []
            images = []

            # Extract text from slides
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)

                    # Extract images from shapes
                    if shape.shape_type == 13:  # Placeholder for image type
                        image = shape.image.blob
                        images.append(image)

            return "\n".join(text), images
        except Exception as e:
            raise ValueError(f"Failed to extract text and images from PPTX: {e}")


class ODPExtractor(PresentationExtractor):
    """
    Extracts text and images from ODP files.
    """

    def extract(self, buffer: bytes) -> Tuple[str, List[bytes]]:
        try:
            # Load the ODP file
            odf_doc = odf_load(io.BytesIO(buffer))
            text = []
            images = []

            # Extract text
            for paragraph in odf_doc.getElementsByType(P):
                text.append(paragraph.textContent)

            # Extract images
            for frame in odf_doc.getElementsByType(Frame):
                for image in frame.getElementsByType(ODFImage):
                    href = image.getAttribute("href")
                    if href:
                        images.append(href)

            return "\n".join(text), images
        except Exception as e:
            raise ValueError(f"Failed to extract text and images from ODP: {e}")


class KeynoteExtractor(PresentationExtractor):
    """
    Extracts text and images from Apple Keynote files.
    """

    def extract(self, buffer: bytes) -> Tuple[str, List[bytes]]:
        try:
            with ZipFile(io.BytesIO(buffer)) as keynote_zip:
                text = []
                images = []

                for file_name in keynote_zip.namelist():
                    if file_name.endswith(".xml"):
                        # Extract text from XML content
                        xml_content = keynote_zip.read(file_name).decode(
                            "utf-8", errors="replace"
                        )
                        text.append(xml_content)

                    # Extract images
                    if file_name.startswith("Data/") and file_name.lower().endswith(
                        (".png", ".jpg", ".jpeg", ".gif")
                    ):
                        images.append(keynote_zip.read(file_name))

                return "\n".join(text), images
        except Exception as e:
            raise ValueError(f"Failed to extract text and images from Keynote: {e}")
