from abc import ABC, abstractmethod
from typing import Tuple, List
from pptx import Presentation
import io


class PresentationExtractor(ABC):
    """
    Abstract base class for slideshow extraction.
    """

    @property
    def file_type(self) -> str:
        return "Slideshow"

    @abstractmethod
    def extract(self, buffer: bytes) -> Tuple[str, List[bytes]]:
        """
        Abstract method to extract text and images from a slideshow.

        Args:
            buffer (bytes): The file content.

        Returns:
            Tuple[str, List[bytes]]: Extracted text and images.
        """
        pass


class PPTXExtractor(PresentationExtractor):
    """
    Extracts text and images from PPTX files.
    """

    def extract(self, buffer: bytes) -> Tuple[str, List[bytes]]:
        try:
            # Open the PPTX file
            presentation = Presentation(io.BytesIO(buffer))
            text = []
            images = []

            # Extract text from slides
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)

                    # Extract images from shapes
                    if shape.shape_type == 13:  # Placeholder for image type
                        image = shape.image.blob
                        images.append(image)

            return "\n".join(text), images
        except Exception as e:
            raise ValueError(f"Failed to extract text and images from PPTX: {e}")